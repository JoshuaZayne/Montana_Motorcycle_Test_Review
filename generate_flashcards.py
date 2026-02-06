"""
Montana Motorcycle Test - Flashcard Generator
==============================================
Generates study flashcards from the Montana Motorcycle Supplement in:
  1. PDF   - Printable two-sided flashcards
  2. PPTX  - PowerPoint slideshow (question -> answer)
  3. HTML  - Interactive browser-based quiz
  4. CSV   - For import into Anki / Quizlet / other apps

Run:  python generate_flashcards.py
Output goes into the same directory as this script (or FLASHCARD_OUTPUT_DIR if set).

Docker:
  docker build -t mt-moto-flashcards .
  docker run -v ./output:/app/output mt-moto-flashcards
"""

import csv
import os
import textwrap
from pathlib import Path

# ── Third-party ──────────────────────────────────────────────────────────────
from reportlab.lib.pagesizes import letter
from reportlab.lib.units import inch
from reportlab.lib.colors import HexColor
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
    PageBreak, Frame, PageTemplate
)
from reportlab.lib.enums import TA_CENTER, TA_LEFT

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Output directory ─────────────────────────────────────────────────────────
OUT_DIR = Path(os.environ.get("FLASHCARD_OUTPUT_DIR", Path(__file__).resolve().parent))
PDF_FILE  = OUT_DIR / "MT_Motorcycle_Flashcards.pdf"
PPTX_FILE = OUT_DIR / "MT_Motorcycle_Flashcards.pptx"
HTML_FILE = OUT_DIR / "MT_Motorcycle_Flashcards.html"
CSV_FILE  = OUT_DIR / "MT_Motorcycle_Flashcards.csv"

# ═════════════════════════════════════════════════════════════════════════════
# FLASHCARD DATA  – extracted from the Montana Motorcycle Supplement (12/2015)
# Each card is (category, question, answer)
# ═════════════════════════════════════════════════════════════════════════════
FLASHCARDS = [
    # ── Endorsement & Licensing ──────────────────────────────────────────────
    ("Licensing",
     "What do you need to legally operate a motorcycle on Montana public roads?",
     "A valid driver license with a motorcycle endorsement."),

    ("Licensing",
     "What is the minimum age to obtain a motorcycle endorsement in Montana?",
     "16 years old (or 15 with an approved driver's education course)."),

    ("Licensing",
     "What are the five parts of the Montana motorcycle driver examination?",
     "1) Driving record review, 2) Physical aptitude review, 3) Vision test, "
     "4) Written test, 5) Road test."),

    ("Licensing",
     "How long is a Montana motorcycle learner permit valid?",
     "One year from the date you pass the knowledge test."),

    ("Licensing",
     "What does a motorcycle learner permit allow you to do?",
     "Operate any motor-driven cycle under the immediate supervision of a "
     "licensed driver who has a valid license for motor-driven cycles."),

    ("Licensing",
     "What is the annual fee for a motorcycle endorsement in Montana?",
     "$0.50 cents per year."),

    # ── Equipment Requirements ───────────────────────────────────────────────
    ("Equipment",
     "How many headlamps must your motorcycle have, and when must they be on?",
     "At least 1 (but not more than 2), illuminated at ALL times. "
     "Must illuminate at least 500 feet ahead."),

    ("Equipment",
     "How far must your taillight and reflector be visible?",
     "At least 500 feet."),

    ("Equipment",
     "How far must your stoplight be visible in normal sunlight?",
     "100 feet."),

    ("Equipment",
     "What is the rearview mirror requirement?",
     "It must provide a clear view of the road for at least 200 feet to the rear."),

    ("Equipment",
     "What brake requirement must your motorcycle meet?",
     "At least one brake that can be operated by foot or hand."),

    ("Equipment",
     "What muffler rules apply to motorcycles in Montana?",
     "Your bike must have a muffler with no modifications that increase sound. "
     "Riding in a national park requires a spark arrestor."),

    ("Equipment",
     "What additional equipment should a motorcycle have beyond the legal minimum?",
     "Passenger footrests/pegs, a horn audible for 200 feet, and a license plate light."),

    # ── Safety Awareness (Emergency Vehicles, Speed) ─────────────────────────
    ("Safety Awareness",
     "What must you do when an emergency vehicle uses audio and visual signals?",
     "Yield the right-of-way, immediately drive to a position parallel and as close "
     "as possible to the right-hand edge or curb of the road."),

    ("Safety Awareness",
     "On highways 50+ MPH, what should you do near a stationary emergency vehicle?",
     "Slow to at least 20 MPH below the posted limit and, when possible, move to a "
     "lane farther away from the stationary vehicle."),

    # ── Protective Gear ──────────────────────────────────────────────────────
    ("Protective Gear",
     "What are the three most important items of protective gear?",
     "1) A DOT-compliant helmet, 2) Face or eye protection, 3) Protective clothing."),

    ("Protective Gear",
     "What are the two primary types of helmets?",
     "Three-quarter (or open-face) and full-face."),

    ("Protective Gear",
     "How should a helmet fit properly?",
     "It should meet U.S. DOT and state standards, fit snugly all the way around, "
     "and have no obvious defects (loose padding, frayed straps, cracks)."),

    ("Protective Gear",
     "What statistics make helmets critical?",
     "One out of every five motorcycle crashes results in head/neck injuries. "
     "Helmeted riders are 3x more likely to survive head injuries. "
     "Most crashes happen < 5 miles from home at < 30 mph."),

    ("Protective Gear",
     "What must effective eye or face protection do?",
     "Be free of scratches, be resistant to penetration, give a clear view to "
     "either side, fasten securely, permit airflow to reduce fogging, and permit "
     "room for eyeglasses if needed."),

    ("Protective Gear",
     "Why shouldn't you rely on a windshield for eye protection?",
     "Most windshields will not protect your eyes from the wind, and neither will "
     "eyeglasses or sunglasses. They might blow off when you turn your head."),

    ("Protective Gear",
     "What clothing is recommended for motorcycle riding?",
     "Jacket and pants covering arms/legs completely (leather is best), sturdy "
     "boots/shoes with slip-resistant soles covering ankles, and leather gloves."),

    ("Protective Gear",
     "Why is hearing protection recommended?",
     "Long-term exposure to engine and wind noise can cause permanent hearing damage. "
     "Hearing protection still lets you hear important sounds like horns and sirens."),

    # ── Know Your Motorcycle ─────────────────────────────────────────────────
    ("Know Your Motorcycle",
     "What should your street-legal motorcycle have at minimum?",
     "Headlight, taillight, brakelight, front and rear brakes, turn signals, "
     "horn, and two mirrors."),

    ("Know Your Motorcycle",
     "How should you choose the right motorcycle for you?",
     "It should 'fit' you: feet reach the ground when seated, controls are easy to "
     "operate. Smaller motorcycles are usually easier for beginners."),

    ("Know Your Motorcycle",
     "What should you do before riding a borrowed or unfamiliar motorcycle?",
     "Make all the checks you would on your own motorcycle, find out where "
     "everything is (turn signals, horn, headlight switch, fuel-supply valve, "
     "engine cut-off switch), and ride cautiously in a controlled area first."),

    # ── T-CLOCS Pre-Ride Inspection ──────────────────────────────────────────
    ("T-CLOCS",
     "What does T-CLOCS stand for?",
     "T = Tires & Wheels, C = Controls, L = Lights & Electrics, "
     "O = Oil & Other Fluids, C = Chassis, S = Stands."),

    ("T-CLOCS",
     "What should you check for Tires & Wheels (T)?",
     "Tire inflation pressure, treadwear, condition of sidewalls, and tread surface."),

    ("T-CLOCS",
     "What should you check for Controls (C)?",
     "Clutch and throttle operate smoothly; throttle snaps back to closed when released; "
     "both brakes feel firm and hold the motorcycle."),

    ("T-CLOCS",
     "What should you check for Lights & Electrics (L)?",
     "Headlight (high and low beams), taillight, both turn signals, "
     "both brakes activate the brake light."),

    ("T-CLOCS",
     "What should you check for Oil & Other Fluids (O)?",
     "Engine oil and transmission fluid levels, brake hydraulic fluid, "
     "coolant level weekly, and make sure the fuel valve is open."),

    ("T-CLOCS",
     "What should you check for Chassis (C)?",
     "Front suspension (no binding), rear shocks move smoothly, "
     "chain adjusted per specs, sprockets not worn or damaged."),

    ("T-CLOCS",
     "What should you check for Stands (S)?",
     "Side stand operates smoothly with spring holding it tightly up; "
     "center stand (if equipped) held firmly against the frame when moving."),

    # ── Body Position & Basic Control ────────────────────────────────────────
    ("Basic Control",
     "What is the correct riding posture?",
     "Sit far enough forward so arms are slightly bent at handgrips. "
     "Keep wrists flat, knees against the tank, feet firmly on footrests "
     "with toes not pointing downward."),

    ("Basic Control",
     "Why should you keep your knees against the gas tank?",
     "It helps you keep your balance as the motorcycle turns."),

    ("Basic Control",
     "What is the typical motorcycle gear pattern?",
     "1-N-2-3-4-5. Neutral is a 'half lift' from 1st gear or 'half press' from 2nd."),

    ("Basic Control",
     "What is the 3-step process for upshifting?",
     "1) Roll off the throttle as you squeeze the clutch lever, "
     "2) Lift the shift lever firmly, "
     "3) Smoothly ease out the clutch and adjust the throttle."),

    ("Basic Control",
     "What is the 3-step process for downshifting?",
     "1) Roll off the throttle as you squeeze the clutch lever, "
     "2) Press the shift lever down firmly, "
     "3) Ease out the clutch as you roll on the throttle."),

    ("Basic Control",
     "Why is engine braking useful?",
     "Shifting to a lower gear causes an effect similar to using the brakes, "
     "helping you slow down. It also helps on downhills and before turns."),

    # ── Braking ──────────────────────────────────────────────────────────────
    ("Braking",
     "How much of total stopping power does the front brake provide?",
     "About 70% of total stopping power."),

    ("Braking",
     "What is the correct braking technique for a straight-line stop?",
     "Use BOTH brakes every time. Squeeze the front brake and apply "
     "light-to-lighter pressure on the rear brake pedal."),

    ("Braking",
     "What happens if the front wheel locks?",
     "Immediate loss of steering control and balance. Release the front brake "
     "immediately and completely, then reapply smoothly."),

    ("Braking",
     "What happens if the rear wheel locks?",
     "The back of the motorcycle may jerk or sway side to side. Keep the rear "
     "brake locked and steer straight until stopped if the wheels are aligned. "
     "If out of alignment, releasing can cause a high-side crash."),

    ("Braking",
     "How does braking in a curve differ from straight-line braking?",
     "Less traction is available because some is used for cornering. "
     "If possible, straighten the motorcycle first, then brake. If you must "
     "brake in a curve, do so smoothly and gradually."),

    ("Braking",
     "What is ABS designed to do?",
     "Prevent wheel lock-up and avoid skids during straight-line panic stops. "
     "It releases and reapplies brake pressure more than 15 times per second."),

    ("Braking",
     "What is the best way to stop quickly?",
     "Apply both brakes at the same time. Squeeze the front brake firmly "
     "and progressively; press down on the rear brake."),

    # ── Turning ──────────────────────────────────────────────────────────────
    ("Turning",
     "What are the four steps for better turning control?",
     "SLOW - reduce speed before the turn. "
     "LOOK - look through the turn to where you want to go. "
     "PRESS - press the handgrip in the direction of the turn (countersteering). "
     "ROLL - roll on the throttle to stabilize the motorcycle."),

    ("Turning",
     "How does countersteering work?",
     "Press left handgrip = lean left = go left. "
     "Press right handgrip = lean right = go right. "
     "The higher the speed, the greater the lean angle needed."),

    ("Turning",
     "In normal turns, how should the rider and motorcycle lean?",
     "Both the rider and motorcycle should lean together at the same angle."),

    ("Turning",
     "How do slow, tight turns differ from normal turns?",
     "In slow, tight turns you counterbalance by leaning the motorcycle only "
     "and keeping your body straight (or leaning opposite)."),

    # ── Lane Position ────────────────────────────────────────────────────────
    ("Lane Position",
     "What are the three lane positions, and which is generally best?",
     "Position 1 (left third), Position 2 (center), Position 3 (right third). "
     "The left third is the most common default. Position 2 (center) is best if "
     "hazards are on both sides."),

    ("Lane Position",
     "Why should you avoid the center of the lane?",
     "The center third is where oil, debris, and grease from cars collect, and where "
     "hazards like manhole covers are found. It can be very slippery when wet."),

    ("Lane Position",
     "What should your lane position help you do?",
     "Increase your ability to see and be seen, avoid others' blind spots, "
     "avoid surface hazards, protect your lane from other drivers, set up for turns, "
     "and communicate your intentions."),

    # ── Following Distance ───────────────────────────────────────────────────
    ("Following Distance",
     "What is the minimum following distance for motorcycles?",
     "A minimum of two seconds behind the vehicle ahead."),

    ("Following Distance",
     "How do you gauge following distance?",
     "Pick out a marker (lamppost, pavement marking). When the vehicle ahead passes "
     "it, count 'one-thousand-one, one-thousand-two.' If you reach the marker before "
     "'two,' you are following too closely."),

    ("Following Distance",
     "When should you increase your following distance beyond 2 seconds?",
     "When pavement is slippery, you cannot see through the vehicle ahead, "
     "traffic is heavy, or someone may squeeze in front of you. "
     "Use a 3-second or greater cushion."),

    # ── Being Followed & Tailgaters ──────────────────────────────────────────
    ("Tailgaters",
     "What is the best way to handle a tailgater?",
     "Change lanes and let them pass. If that's not possible, slow down and open "
     "up extra space ahead of you so both you and the tailgater have room to stop."),

    ("Tailgaters",
     "When should you flash your brake light?",
     "When you slow more quickly than others might expect (e.g., mid-block, alley) "
     "or when a tailgater is following too closely."),

    # ── Passing & Being Passed ───────────────────────────────────────────────
    ("Passing",
     "What are the four steps for passing another vehicle?",
     "1) Ride in the left portion of the lane, signal, and check for traffic. "
     "2) When safe, move into the left lane and accelerate. "
     "3) Ride through the blind spot quickly. "
     "4) Signal again, check mirrors, return to your lane, cancel the signal."),

    ("Passing",
     "What should you do when being passed?",
     "Stay in the center of your lane. Do not move to the far side of the lane; "
     "that invites the other driver to cut back in too early."),

    # ── Lane Sharing ─────────────────────────────────────────────────────────
    ("Lane Sharing",
     "Is lane sharing legal in Montana?",
     "Lane sharing is usually prohibited. Cars and motorcycles need a full lane "
     "to operate safely."),

    ("Lane Sharing",
     "When should you ride center-lane position to discourage lane sharing?",
     "In heavy bumper-to-bumper traffic, when other drivers want to pass you, "
     "when you are preparing to turn at an intersection, or moving into an exit lane."),

    # ── Intersections ────────────────────────────────────────────────────────
    ("Intersections",
     "Where do most motorcycle/car crashes occur?",
     "At intersections. Over half of motorcycle/car crashes are caused by "
     "drivers entering a rider's right-of-way."),

    ("Intersections",
     "Should you rely on 'eye contact' with a driver at an intersection?",
     "No. A driver may look right at you and still fail to 'see' you. "
     "If a car can enter your path, assume it will."),

    ("Intersections",
     "What should you do at a blind intersection?",
     "Move to the lane portion that brings you into another driver's field of "
     "vision at the earliest possible moment."),

    ("Intersections",
     "What should you do at a stop sign?",
     "Stop first at the stop line, then edge forward and stop again just short "
     "of where the cross-traffic lane meets your lane. Lean forward and look "
     "around buildings, parked cars, or bushes. Keep your front wheel out of "
     "the cross lane of travel."),

    # ── Increasing Conspicuity (Visibility) ──────────────────────────────────
    ("Visibility",
     "What colors are most visible to other drivers?",
     "Bright orange, red, yellow, or green jackets/vests. Reflective, "
     "bright-colored helmet and clothing are best."),

    ("Visibility",
     "Why should your headlight always be on?",
     "A motorcycle with its headlight on during the day is twice as likely to be "
     "noticed. Use low beam during the day and in fog."),

    ("Visibility",
     "What should you do with your signal when entering a freeway?",
     "Keep it blinking; drivers approaching from behind are more likely to see "
     "it and make room for you."),

    # ── SEE Strategy ─────────────────────────────────────────────────────────
    ("SEE Strategy",
     "What does SEE stand for?",
     "Search, Evaluate, Execute - a three-step strategy to reduce risk."),

    ("SEE Strategy",
     "What are the three 'lead times' experienced riders use when Searching?",
     "2 seconds - be alert and scan for immediate hazards. "
     "4 seconds - quick response zone for things going wrong. "
     "12 seconds - look far ahead for developing hazards to prepare early."),

    ("SEE Strategy",
     "What does 'Evaluate' mean in the SEE strategy?",
     "Think about how hazards can interact to create risks. Anticipate potential "
     "problems and have a plan, especially for road/surface hazards, traffic "
     "control devices, and other vehicles."),

    ("SEE Strategy",
     "What does 'Execute' mean in the SEE strategy?",
     "Create more space and minimize harm by: communicating (lights, horn), "
     "adjusting speed (accelerate, stop, slow), adjusting position (swerve, "
     "change lanes, move within lane)."),

    ("SEE Strategy",
     "To reduce your reaction time, you should:",
     "Cover the clutch and the brakes so you can react faster when a hazard appears."),

    # ── Crash Avoidance ──────────────────────────────────────────────────────
    ("Crash Avoidance",
     "What two skills are critical for crash avoidance?",
     "Knowing when to stop or swerve, and being able to do both quickly and properly."),

    ("Crash Avoidance",
     "What do studies show about most crash-involved riders?",
     "They under-braked the front tire and over-braked the rear. They also "
     "did not separate braking from swerving when appropriate."),

    ("Crash Avoidance",
     "How should you swerve to avoid an obstacle?",
     "Press the handgrip in the direction of the turn (countersteering). "
     "Do NOT brake while swerving - either brake then swerve, or swerve then brake."),

    ("Crash Avoidance",
     "What is the correct technique for stopping quickly in a curve?",
     "If possible, straighten the motorcycle first, then brake. If you can't "
     "straighten, apply brakes smoothly and gradually, reducing lean angle as you slow."),

    # ── Cornering ────────────────────────────────────────────────────────────
    ("Cornering",
     "What is a primary cause of single-vehicle motorcycle crashes?",
     "Riders running wide in a curve or turn and hitting a fixed object or "
     "going off the road."),

    ("Cornering",
     "How should you handle decreasing-radius (tighter) curves?",
     "Start at the outside of the curve to increase your line of sight, move "
     "toward the inside as you turn, then move to the outside to exit."),

    # ── Handling Dangerous Surfaces ──────────────────────────────────────────
    ("Dangerous Surfaces",
     "Name five surfaces that provide poor traction for motorcycles.",
     "1) Wet pavement, 2) Gravel roads, 3) Mud/leaves/snow/ice, "
     "4) Painted lane markings (wet), 5) Steel plates and manhole covers."),

    ("Dangerous Surfaces",
     "How should you ride on a slippery surface?",
     "Reduce speed, avoid sudden moves, use both brakes gently, squeeze "
     "(don't grab) the front brake to avoid locking."),

    ("Dangerous Surfaces",
     "When it starts to rain, which part of the lane is safest?",
     "The tire tracks left by cars. The center of the lane is most hazardous "
     "when wet because that's where oil and fluids collect."),

    ("Dangerous Surfaces",
     "How should you cross railroad tracks?",
     "Ride straight across (head-on at a 90-degree angle) whenever possible. "
     "If tracks are parallel, cross at a minimum 45-degree angle. "
     "Edging across can catch your tires and throw you off balance."),

    ("Dangerous Surfaces",
     "How do you ride over an obstacle you cannot avoid?",
     "Slow down as much as possible, make sure the motorcycle is straight, "
     "rise slightly off the seat on the footrests, roll on the throttle slightly "
     "just before contact to lighten the front end."),

    # ── Riding at Night ──────────────────────────────────────────────────────
    ("Night Riding",
     "What three things should you do when riding at night?",
     "1) Reduce your speed below daytime levels, 2) Increase your following "
     "distance to at least 3 seconds, 3) Use the car ahead's headlights to "
     "preview the road and use your high beam when not following or meeting traffic."),

    ("Night Riding",
     "Why are distances harder to judge at night?",
     "Your eyes rely on shadows and light contrasts that may be missing or "
     "distorted under artificial lights, making objects seem farther or closer "
     "than they really are."),

    # ── Mechanical Problems ──────────────────────────────────────────────────
    ("Mechanical Problems",
     "What should you do if you get a flat tire while riding?",
     "Hold handgrips firmly, ease off the throttle, and keep a straight course. "
     "If braking is required, gradually apply the brake of the tire that is NOT "
     "flat. Edge to the side of the road and stop."),

    ("Mechanical Problems",
     "What should you do if the throttle is stuck?",
     "Twist it back and forth to free it. If that doesn't work, immediately "
     "operate the engine cut-off switch and pull in the clutch at the same time."),

    ("Mechanical Problems",
     "How should you handle a 'wobble' (tank-slapper)?",
     "Grip the handlebars firmly but don't fight the wobble. Close the throttle "
     "gradually (don't brake). Move your weight as far forward and down as possible. "
     "Pull off the road as soon as you can."),

    ("Mechanical Problems",
     "What causes a wobble?",
     "Improper loading, unsuitable accessories, incorrect tire pressure, "
     "worn steering parts, misaligned or out-of-balance front wheel, "
     "loose wheel bearings or spokes, or worn swingarm bearings."),

    ("Mechanical Problems",
     "What should you do if the chain or belt breaks?",
     "You will notice an instant loss of power. Close the throttle and brake "
     "to a stop in a safe area. Loss of oil in the rear differential can "
     "cause the rear wheel to lock."),

    ("Mechanical Problems",
     "What should you do if the engine seizes ('locks up')?",
     "Squeeze the clutch to disengage the engine from the rear wheel. "
     "Pull off the road and stop. Check the oil; add oil if needed, otherwise "
     "let the engine cool before restarting."),

    # ── Animals ───────────────────────────────────────────────────────────────
    ("Animals",
     "What should you do if an animal enters your lane?",
     "Hitting something small is less dangerous than hitting something big "
     "(like a car). Stay in your lane. For larger animals (deer, elk, cattle), "
     "brake and prepare to stop - they are unpredictable."),

    ("Animals",
     "What should you do if a dog chases you?",
     "Approach the animal slowly, then accelerate away and leave it behind. "
     "Do NOT kick at it."),

    # ── Flying Objects ───────────────────────────────────────────────────────
    ("Flying Objects",
     "How should you handle being hit by a flying object (insects, debris)?",
     "Keep your eyes on the road and your hands on the handlebars. "
     "When safe, pull off the road and repair any damage. Face protection "
     "is essential to prevent objects from hitting your eyes or face."),

    # ── Getting Off the Road ─────────────────────────────────────────────────
    ("Getting Off Road",
     "What three things should you do when leaving the road?",
     "1) Check the roadside - make sure it's firm enough to ride on. "
     "2) Signal your intentions to drivers behind you. "
     "3) Pull off as far as you can. It's hard for others to spot a motorcycle "
     "by the side of the road."),

    # ── Carrying Passengers ──────────────────────────────────────────────────
    ("Passengers",
     "What equipment must your motorcycle have to carry a passenger?",
     "Passenger footrests/pegs, a proper seat large enough for two without "
     "crowding, and a strap or solid handholds for the passenger."),

    ("Passengers",
     "What instructions should you give a passenger before riding?",
     "Mount only after the engine is started and transmission is in neutral; "
     "sit as far forward as possible; hold your waist, hips, belt, or passenger "
     "handholds firmly; keep both feet on the footrests; lean with you; "
     "avoid unnecessary conversation and sudden moves; keep legs away from "
     "muffler(s), chains, and moving parts."),

    ("Passengers",
     "How does carrying a passenger affect riding?",
     "The motorcycle responds more slowly - it takes longer to speed up, "
     "slow down, or turn. Ride a little slower in curves, start slowing earlier, "
     "and maintain a larger space cushion."),

    ("Passengers",
     "Where should a child passenger be placed?",
     "Children should be placed immediately behind the rider. A child sitting "
     "in front will not be able to properly balance and may interfere with "
     "the rider's control."),

    # ── Carrying Cargo ───────────────────────────────────────────────────────
    ("Cargo",
     "How should you load cargo on a motorcycle?",
     "Keep the load forward and low. Pack heavier items in the front of the "
     "tank bag. If using saddlebags, load each with about the same weight. "
     "Fasten the load securely with elastic cords (one attachment point per side "
     "recommended). Do not exceed gross vehicle weight rating."),

    ("Cargo",
     "What can happen if cargo is improperly loaded?",
     "An uneven load can cause the motorcycle to pull to one side. Overloading "
     "can cause bags to catch in the wheel or chain, locking the rear wheel and "
     "causing a skid."),

    # ── Group Riding ─────────────────────────────────────────────────────────
    ("Group Riding",
     "What formation should a group of motorcyclists use?",
     "A staggered formation. The leader rides in the left third of the lane; "
     "the second rider stays at least one second behind in the right third; "
     "the third rider is at least two seconds behind the leader in the left "
     "third, and so on."),

    ("Group Riding",
     "When should a group switch from staggered to single-file formation?",
     "When riding in curves, turning, entering or leaving freeways/highways, "
     "and in any situation where the full lane width is needed."),

    ("Group Riding",
     "What are the most important rules for group riding?",
     "No competition, no passing of other riders, and no tailgating. "
     "Place inexperienced riders just behind the leader."),

    ("Group Riding",
     "Where should inexperienced riders position themselves in a group?",
     "Just behind the leader so they can keep pace without riding faster "
     "than they are comfortable with, and more experienced riders can watch "
     "them from behind."),

    ("Group Riding",
     "What should happen if a rider falls behind during a group ride?",
     "Everyone should slow down to keep the group together. The lead rider "
     "sets the pace based on the least experienced rider."),

    # ── Alcohol & Drugs ──────────────────────────────────────────────────────
    ("Alcohol & Drugs",
     "What BAC level is considered legally intoxicated in Montana?",
     "0.08% for adults (21+). For operators under 21, lower limits (0.00-0.02%) "
     "apply depending on the state."),

    ("Alcohol & Drugs",
     "How quickly can alcohol be eliminated from the body?",
     "About one drink per hour. Alcohol may still accumulate even if you drink "
     "at a rate of one drink per hour."),

    ("Alcohol & Drugs",
     "What percentage of fatal motorcycle crashes involve alcohol?",
     "Nearly 40% of all riders killed in motorcycle crashes had been drinking."),

    ("Alcohol & Drugs",
     "What are the consequences of a DUI conviction in Montana?",
     "License suspension, severe fines, community service (litter pickup, "
     "washing cars at emergency wards), additional lawyer fees, lost work time, "
     "and the psychological costs of being tagged a 'drunk driver.'"),

    ("Alcohol & Drugs",
     "If you wait one hour per drink before riding, are you safe?",
     "Not necessarily. Side effects from the drinking may still remain and "
     "affect your riding skills. It depends on many factors including body weight, "
     "food intake, and physical condition."),

    # ── Fatigue ───────────────────────────────────────────────────────────────
    ("Fatigue",
     "How does fatigue affect motorcycle riding?",
     "Fatigue can affect your control of the motorcycle. Motorcycling is more "
     "tiring than driving a car, and you tire sooner."),

    ("Fatigue",
     "What are four tips to combat fatigue on long rides?",
     "1) Protect yourself from wind, cold, and rain. "
     "2) Limit your distance (experienced riders suggest ~6 hours/day). "
     "3) Take frequent rest breaks every two hours. "
     "4) Don't drink or use drugs."),

    # ── Using Your Horn ──────────────────────────────────────────────────────
    ("Horn",
     "When should you use your horn?",
     "When a driver in the next lane is close and may want to pass; when a "
     "parked car has someone in the driver's seat; when someone is in the "
     "street riding a bicycle or walking. In an emergency, sound it loud and long."),

    ("Horn",
     "Should you rely on your horn to keep you safe?",
     "No. A motorcycle's horn isn't as loud as a car's. Use it, but also have "
     "time and space to maneuver as backup strategies."),

    # ── Using Mirrors / Head Checks ──────────────────────────────────────────
    ("Mirrors",
     "When should you check your mirrors?",
     "When stopped at an intersection (watch for cars coming up behind), "
     "before you change lanes, before you slow down, and as part of your "
     "normal scanning routine."),

    ("Mirrors",
     "Why are head checks necessary in addition to mirrors?",
     "Motorcycles have blind spots just like cars. Before changing lanes, "
     "turn your head and look to the side for other vehicles."),

    ("Mirrors",
     "Are convex mirrors (rounded) common on motorcycles? What's the caution?",
     "Yes, most motorcycles have convex mirrors. They provide a wider view but "
     "make cars seem farther away than they actually are."),

    # ── Parking ───────────────────────────────────────────────────────────────
    ("Parking",
     "How should you park a motorcycle at a curb (parallel parking)?",
     "Position the motorcycle at an angle with the rear wheel to the curb. "
     "Some cities require motorcycles to park parallel to the curb."),

    # ── Passing Parked Cars ──────────────────────────────────────────────────
    ("Parked Cars",
     "How should you ride when passing parked cars?",
     "Stay toward the left of your lane to avoid problems from doors opening, "
     "drivers getting out, or people stepping between cars. If oncoming traffic "
     "is present, stay in the center-lane position."),

    # ── Hand Signals ─────────────────────────────────────────────────────────
    ("Hand Signals",
     "What is the hand signal for 'Single File'?",
     "Arm raised, index finger extended straight up."),

    ("Hand Signals",
     "What is the hand signal for 'Stop'?",
     "Arm extended straight down, palm facing back."),

    ("Hand Signals",
     "What is the hand signal for 'Speed Up'?",
     "Arm extended straight out, palm facing up, swing upward."),

    ("Hand Signals",
     "What is the hand signal for 'Slow Down'?",
     "Arm extended straight out, palm facing down, swing down to your side."),

    ("Hand Signals",
     "What is the hand signal for 'Follow Me'?",
     "Arm extended straight up from shoulder, palm forward."),

    ("Hand Signals",
     "What is the hand signal for 'Hazard in Roadway'?",
     "On the left - point with left hand. On the right - point with right foot."),

    ("Hand Signals",
     "What is the hand signal for 'Fuel'?",
     "Arm out to side pointing to tank with finger extended."),

    ("Hand Signals",
     "What is the hand signal for 'Comfort Stop'?",
     "Forearm extended, fist clenched, with short up and down motion."),

    ("Hand Signals",
     "What is the hand signal for 'Refreshment Stop'?",
     "Fingers closed, thumb to mouth."),

    # ── Test Yourself Questions (from the manual) ────────────────────────────
    ("Practice Test",
     "TEST Q1: A plastic shatter-resistant face shield: (A) is not necessary if "
     "you have a windshield (B) only protects your eyes (C) helps protect your "
     "whole face (D) does not protect your face as well as goggles",
     "C - Helps protect your whole face."),

    ("Practice Test",
     "TEST Q2: More than half of all crashes: (A) occur at speeds > 35 mph "
     "(B) happen at night (C) are caused by worn tires "
     "(D) involve riders with < 5 months experience on their motorcycles",
     "D - Involve riders who have less than five months of experience on "
     "their motorcycles."),

    ("Practice Test",
     "TEST Q3: When riding, you should: (A) turn your head and shoulders to "
     "look through turns (B) keep your arms straight (C) keep your knees away "
     "from the gas tank (D) turn just your head and eyes to look where you are going",
     "D - Turn just your head and eyes to look where you are going."),

    ("Practice Test",
     "TEST Q4: Usually, a good way to handle tailgaters is to:",
     "A - Change lanes and let them pass."),

    ("Practice Test",
     "TEST Q5: To reduce your reaction time, you should:",
     "B - Cover the clutch and the brakes."),

    ("Practice Test",
     "TEST Q6: Making eye contact with other drivers:",
     "C - Doesn't mean that the driver will yield."),

    ("Practice Test",
     "TEST Q7: Reflective clothing should be worn:",
     "D - Be worn day and night."),

    ("Practice Test",
     "TEST Q8: The best way to stop quickly is to:",
     "D - Use both brakes at the same time."),

    ("Practice Test",
     "TEST Q9: When it starts to rain it is usually best to:",
     "C - Ride in the tire tracks left by cars."),

    ("Practice Test",
     "TEST Q10: If your motorcycle starts to wobble:",
     "C - Grip the handlebars firmly and close the throttle gradually."),

    ("Practice Test",
     "TEST Q11: If you are chased by a dog:",
     "D - Approach the animal slowly, then speed up."),

    ("Practice Test",
     "TEST Q12: Passengers should:",
     "A - Lean as you lean."),

    ("Practice Test",
     "TEST Q13: When riding in a group, inexperienced riders should position themselves:",
     "A - Just behind the leader."),

    ("Practice Test",
     "TEST Q14: If you wait one hour per drink for alcohol to be eliminated before riding:",
     "C - Side effects from the drinking may still remain."),

    # ── Knowledge Test Sample Questions (from page 46-47) ────────────────────
    ("Knowledge Test",
     "SAMPLE Q1: It is MOST important to flash your brake light when:",
     "B - Someone is following too closely."),

    ("Knowledge Test",
     "SAMPLE Q2: The FRONT brake supplies how much of the potential stopping power?",
     "C - About 70%."),

    ("Knowledge Test",
     "SAMPLE Q3: To swerve correctly:",
     "C - Press the handgrip in the direction of the turn."),

    ("Knowledge Test",
     "SAMPLE Q4: If a tire goes flat while riding and you must stop, it is usually best to:",
     "C - Brake on the good tire and steer to the side of the road."),

    ("Knowledge Test",
     "SAMPLE Q5: The car below is waiting to enter the intersection. It is best to:",
     "B - Reduce speed and be ready to react."),

    # ── Skill Test Topics ────────────────────────────────────────────────────
    ("Skill Test",
     "What abilities are tested in the on-motorcycle skill test?",
     "Know your motorcycle and riding limits; accelerate, brake, and turn safely; "
     "see, be seen, and communicate; adjust speed and position to traffic; "
     "stop, turn, and swerve quickly; make critical decisions and carry them out."),

    ("Skill Test",
     "What factors do examiners score during the skill test?",
     "Selecting safe speeds, choosing correct path and staying within boundaries, "
     "completing normal and quick stops, completing normal and quick turns or swerves."),
]


# ═════════════════════════════════════════════════════════════════════════════
#  1. CSV  OUTPUT
# ═════════════════════════════════════════════════════════════════════════════
def generate_csv():
    with open(CSV_FILE, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(["Category", "Question", "Answer"])
        for cat, q, a in FLASHCARDS:
            writer.writerow([cat, q, a])
    print(f"  CSV  -> {CSV_FILE}")


# ═════════════════════════════════════════════════════════════════════════════
#  2. PDF  OUTPUT  (printable flashcards)
# ═════════════════════════════════════════════════════════════════════════════
def generate_pdf():
    doc = SimpleDocTemplate(
        str(PDF_FILE),
        pagesize=letter,
        topMargin=0.5 * inch,
        bottomMargin=0.5 * inch,
        leftMargin=0.6 * inch,
        rightMargin=0.6 * inch,
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "CardTitle",
        parent=styles["Heading1"],
        fontSize=22,
        alignment=TA_CENTER,
        spaceAfter=20,
        textColor=HexColor("#1a1a2e"),
    )
    cat_style = ParagraphStyle(
        "Category",
        parent=styles["Heading3"],
        fontSize=11,
        alignment=TA_CENTER,
        textColor=HexColor("#e94560"),
        spaceAfter=6,
    )
    q_style = ParagraphStyle(
        "Question",
        parent=styles["Normal"],
        fontSize=13,
        alignment=TA_CENTER,
        leading=18,
        textColor=HexColor("#16213e"),
        spaceAfter=10,
        spaceBefore=10,
    )
    a_label = ParagraphStyle(
        "AnswerLabel",
        parent=styles["Heading3"],
        fontSize=11,
        alignment=TA_CENTER,
        textColor=HexColor("#0f3460"),
        spaceAfter=4,
    )
    a_style = ParagraphStyle(
        "Answer",
        parent=styles["Normal"],
        fontSize=12,
        alignment=TA_CENTER,
        leading=17,
        textColor=HexColor("#1a1a2e"),
        spaceAfter=10,
        spaceBefore=4,
    )
    footer_style = ParagraphStyle(
        "Footer",
        parent=styles["Normal"],
        fontSize=8,
        alignment=TA_CENTER,
        textColor=HexColor("#999999"),
    )

    story = []
    # Cover page
    story.append(Spacer(1, 1.5 * inch))
    story.append(Paragraph("Montana Motorcycle Test", title_style))
    story.append(Paragraph("Flashcard Study Guide", ParagraphStyle(
        "Sub", parent=styles["Heading2"], alignment=TA_CENTER,
        fontSize=16, textColor=HexColor("#0f3460"), spaceAfter=30)))
    story.append(Paragraph(
        f"{len(FLASHCARDS)} Flashcards Covering All Manual Topics",
        ParagraphStyle("Count", parent=styles["Normal"],
                       alignment=TA_CENTER, fontSize=12,
                       textColor=HexColor("#555555"), spaceAfter=20)))
    story.append(Spacer(1, 0.5 * inch))

    categories = sorted(set(c for c, _, _ in FLASHCARDS))
    toc_lines = "<br/>".join(
        f"&bull; {cat} ({sum(1 for c,_,_ in FLASHCARDS if c==cat)} cards)"
        for cat in categories
    )
    story.append(Paragraph(
        f"<b>Categories:</b><br/><br/>{toc_lines}",
        ParagraphStyle("TOC", parent=styles["Normal"],
                       fontSize=11, alignment=TA_CENTER, leading=16,
                       textColor=HexColor("#333333"))))
    story.append(PageBreak())

    # Flashcards - one card per page: question side, then answer side
    for i, (cat, question, answer) in enumerate(FLASHCARDS, 1):
        # ── Question side ──
        story.append(Spacer(1, 0.3 * inch))
        story.append(Paragraph(f"CARD {i} of {len(FLASHCARDS)}", footer_style))
        story.append(Spacer(1, 0.5 * inch))
        story.append(Paragraph(cat.upper(), cat_style))
        story.append(Spacer(1, 0.2 * inch))

        # Draw a box around the question
        q_data = [[Paragraph(f"<b>Q:</b> {question}", q_style)]]
        q_table = Table(q_data, colWidths=[6.5 * inch], rowHeights=None)
        q_table.setStyle(TableStyle([
            ("BOX", (0, 0), (-1, -1), 2, HexColor("#e94560")),
            ("TOPPADDING", (0, 0), (-1, -1), 20),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 20),
            ("LEFTPADDING", (0, 0), (-1, -1), 15),
            ("RIGHTPADDING", (0, 0), (-1, -1), 15),
            ("BACKGROUND", (0, 0), (-1, -1), HexColor("#fef9ff")),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ]))
        story.append(q_table)

        story.append(Spacer(1, 1.0 * inch))
        story.append(Paragraph("( flip for answer )", footer_style))
        story.append(PageBreak())

        # ── Answer side ──
        story.append(Spacer(1, 0.3 * inch))
        story.append(Paragraph(f"CARD {i} of {len(FLASHCARDS)} - ANSWER", footer_style))
        story.append(Spacer(1, 0.5 * inch))
        story.append(Paragraph(cat.upper(), cat_style))
        story.append(Spacer(1, 0.1 * inch))

        # Question recap (smaller)
        story.append(Paragraph(
            f"<i>Q: {question}</i>",
            ParagraphStyle("QRecap", parent=styles["Normal"],
                           fontSize=10, alignment=TA_CENTER,
                           textColor=HexColor("#888888"), leading=14,
                           spaceAfter=12)))

        # Answer box
        a_data = [[Paragraph(f"{answer}", a_style)]]
        a_table = Table(a_data, colWidths=[6.5 * inch], rowHeights=None)
        a_table.setStyle(TableStyle([
            ("BOX", (0, 0), (-1, -1), 2, HexColor("#0f3460")),
            ("TOPPADDING", (0, 0), (-1, -1), 20),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 20),
            ("LEFTPADDING", (0, 0), (-1, -1), 15),
            ("RIGHTPADDING", (0, 0), (-1, -1), 15),
            ("BACKGROUND", (0, 0), (-1, -1), HexColor("#f0f7ff")),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ]))
        story.append(a_table)
        story.append(PageBreak())

    doc.build(story)
    print(f"  PDF  -> {PDF_FILE}")


# ═════════════════════════════════════════════════════════════════════════════
#  3. POWERPOINT  OUTPUT
# ═════════════════════════════════════════════════════════════════════════════
def generate_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Helper to add a colored rectangle ────────────────────────────────────
    def add_rect(slide, left, top, width, height, fill_color):
        from pptx.util import Emu
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    def add_text_box(slide, left, top, width, height, text, font_size=18,
                     bold=False, color=RGBColor(0x1a, 0x1a, 0x2e),
                     alignment=PP_ALIGN.CENTER):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = alignment
        return txBox

    # ── Title slide ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height,
             RGBColor(0x1a, 0x1a, 0x2e))
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11.333), Inches(1.5),
                 "MONTANA MOTORCYCLE TEST", 44, True, RGBColor(0xFF, 0xFF, 0xFF))
    add_text_box(slide, Inches(1), Inches(3.2), Inches(11.333), Inches(1),
                 "Flashcard Study Guide", 28, False, RGBColor(0xe9, 0x45, 0x60))
    add_text_box(slide, Inches(1), Inches(4.5), Inches(11.333), Inches(1),
                 f"{len(FLASHCARDS)} Cards | All Manual Topics | MT MVD Supplement",
                 18, False, RGBColor(0xAA, 0xAA, 0xCC))
    add_text_box(slide, Inches(1), Inches(6), Inches(11.333), Inches(0.8),
                 "Press SPACEBAR or click to advance through cards",
                 14, False, RGBColor(0x88, 0x88, 0xAA))

    # ── Flashcard slides ─────────────────────────────────────────────────────
    for i, (cat, question, answer) in enumerate(FLASHCARDS, 1):
        # QUESTION slide
        q_slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_rect(q_slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2),
                 RGBColor(0xe9, 0x45, 0x60))
        add_text_box(q_slide, Inches(0.5), Inches(0.2), Inches(8), Inches(0.8),
                     f"Card {i}/{len(FLASHCARDS)}  |  {cat.upper()}",
                     16, True, RGBColor(0xFF, 0xFF, 0xFF), PP_ALIGN.LEFT)
        add_text_box(q_slide, Inches(9), Inches(0.2), Inches(4), Inches(0.8),
                     "QUESTION", 20, True, RGBColor(0xFF, 0xFF, 0xFF))

        # Question text centered
        add_text_box(q_slide, Inches(1.5), Inches(2.0), Inches(10.333), Inches(4),
                     question, 24, False, RGBColor(0x1a, 0x1a, 0x2e))

        add_text_box(q_slide, Inches(1), Inches(6.3), Inches(11.333), Inches(0.6),
                     "click or press SPACE for answer >>",
                     12, False, RGBColor(0xAA, 0xAA, 0xAA))

        # ANSWER slide
        a_slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_rect(a_slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2),
                 RGBColor(0x0f, 0x34, 0x60))
        add_text_box(a_slide, Inches(0.5), Inches(0.2), Inches(8), Inches(0.8),
                     f"Card {i}/{len(FLASHCARDS)}  |  {cat.upper()}",
                     16, True, RGBColor(0xFF, 0xFF, 0xFF), PP_ALIGN.LEFT)
        add_text_box(a_slide, Inches(9), Inches(0.2), Inches(4), Inches(0.8),
                     "ANSWER", 20, True, RGBColor(0xFF, 0xFF, 0xFF))

        # Question recap
        add_text_box(a_slide, Inches(1.5), Inches(1.5), Inches(10.333), Inches(1.5),
                     f"Q: {question}", 14, False, RGBColor(0x88, 0x88, 0x88))

        # Answer box background
        add_rect(a_slide, Inches(1.2), Inches(3.2), Inches(10.9), Inches(3.2),
                 RGBColor(0xf0, 0xf7, 0xff))
        add_text_box(a_slide, Inches(1.5), Inches(3.5), Inches(10.333), Inches(2.8),
                     answer, 22, False, RGBColor(0x0f, 0x34, 0x60))

    prs.save(str(PPTX_FILE))
    print(f"  PPTX -> {PPTX_FILE}")


# ═════════════════════════════════════════════════════════════════════════════
#  4. HTML  OUTPUT  (interactive quiz with flip animation)
# ═════════════════════════════════════════════════════════════════════════════
def generate_html():
    import json

    cards_json = json.dumps([
        {"cat": cat, "q": q, "a": a} for cat, q, a in FLASHCARDS
    ], indent=2)

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>Montana Motorcycle Test - Flashcards</title>
<style>
  * {{ margin:0; padding:0; box-sizing:border-box; }}
  body {{
    font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
    background: linear-gradient(135deg, #1a1a2e 0%, #16213e 50%, #0f3460 100%);
    min-height: 100vh; color: #fff;
    display: flex; flex-direction: column; align-items: center;
  }}
  header {{
    text-align: center; padding: 25px 20px 10px;
  }}
  header h1 {{ font-size: 2rem; color: #e94560; }}
  header p {{ color: #aab; font-size: 0.95rem; margin-top: 5px; }}

  .controls {{
    display: flex; gap: 12px; flex-wrap: wrap;
    justify-content: center; padding: 10px 20px 5px;
  }}
  .controls button {{
    padding: 8px 20px; border: none; border-radius: 6px;
    cursor: pointer; font-size: 0.95rem; font-weight: 600;
    transition: transform .1s;
  }}
  .controls button:active {{ transform: scale(.95); }}
  .btn-prev {{ background: #533483; color: #fff; }}
  .btn-flip {{ background: #e94560; color: #fff; }}
  .btn-next {{ background: #533483; color: #fff; }}
  .btn-shuffle {{ background: #0f3460; color: #fff; }}
  .btn-reset {{ background: #444; color: #fff; }}

  .progress {{
    color: #aab; font-size: 0.9rem; padding: 8px;
    text-align: center;
  }}

  .filter-bar {{
    display: flex; flex-wrap: wrap; gap: 6px;
    justify-content: center; padding: 5px 20px 10px;
    max-width: 1000px;
  }}
  .filter-btn {{
    padding: 4px 12px; border: 1px solid #555;
    border-radius: 20px; background: transparent;
    color: #ccc; cursor: pointer; font-size: 0.8rem;
    transition: all .2s;
  }}
  .filter-btn.active {{ background: #e94560; color: #fff; border-color: #e94560; }}

  .card-container {{
    perspective: 1200px; width: 90%; max-width: 700px;
    height: 380px; margin: 10px auto;
  }}
  .card {{
    width: 100%; height: 100%; position: relative;
    transform-style: preserve-3d;
    transition: transform 0.6s cubic-bezier(.4,0,.2,1);
    cursor: pointer;
  }}
  .card.flipped {{ transform: rotateY(180deg); }}
  .card-face {{
    position: absolute; inset: 0; backface-visibility: hidden;
    border-radius: 16px; padding: 30px 35px;
    display: flex; flex-direction: column;
    justify-content: center; align-items: center;
    text-align: center; box-shadow: 0 10px 40px rgba(0,0,0,.4);
  }}
  .card-front {{
    background: linear-gradient(145deg, #fff 0%, #f5f0ff 100%);
    color: #1a1a2e;
  }}
  .card-back {{
    background: linear-gradient(145deg, #f0f7ff 0%, #e8f0fe 100%);
    color: #0f3460;
    transform: rotateY(180deg);
  }}
  .card-cat {{
    font-size: 0.85rem; font-weight: 700; letter-spacing: 1px;
    color: #e94560; text-transform: uppercase; margin-bottom: 15px;
  }}
  .card-text {{ font-size: 1.15rem; line-height: 1.6; }}
  .card-back .card-text {{ font-size: 1.1rem; }}
  .card-label {{
    font-size: 0.75rem; color: #999; margin-top: 15px;
    letter-spacing: 0.5px;
  }}

  .score-bar {{
    display: flex; gap: 10px; justify-content: center;
    padding: 5px 20px 10px;
  }}
  .score-btn {{
    padding: 6px 18px; border: none; border-radius: 6px;
    cursor: pointer; font-size: 0.85rem; font-weight: 600;
  }}
  .score-know {{ background: #27ae60; color: #fff; }}
  .score-unsure {{ background: #f39c12; color: #fff; }}
  .score-dunno {{ background: #c0392b; color: #fff; }}
  .score-display {{
    color: #aab; font-size: 0.85rem; padding: 0 20px 15px;
    text-align: center;
  }}

  .keyboard-hint {{
    color: #667; font-size: 0.8rem; text-align: center;
    padding: 0 20px 20px;
  }}

  @media (max-width: 600px) {{
    .card-container {{ height: 340px; }}
    .card-text {{ font-size: 1rem; }}
    header h1 {{ font-size: 1.4rem; }}
  }}
</style>
</head>
<body>

<header>
  <h1>Montana Motorcycle Test - Flashcards</h1>
  <p>Click the card or press SPACE to flip | Arrow keys to navigate</p>
</header>

<div class="filter-bar" id="filterBar"></div>

<div class="progress" id="progress"></div>

<div class="card-container" id="cardContainer">
  <div class="card" id="card" onclick="flip()">
    <div class="card-face card-front">
      <div class="card-cat" id="qCat"></div>
      <div class="card-text" id="qText"></div>
      <div class="card-label">QUESTION - click to reveal answer</div>
    </div>
    <div class="card-face card-back">
      <div class="card-cat" id="aCat"></div>
      <div class="card-text" id="aText"></div>
      <div class="card-label">ANSWER - click to return to question</div>
    </div>
  </div>
</div>

<div class="controls">
  <button class="btn-prev" onclick="prev()">&#9664; Prev</button>
  <button class="btn-flip" onclick="flip()">Flip Card</button>
  <button class="btn-next" onclick="next()">Next &#9654;</button>
  <button class="btn-shuffle" onclick="shuffleCards()">Shuffle</button>
  <button class="btn-reset" onclick="resetCards()">Reset</button>
</div>

<div class="score-bar">
  <button class="score-know" onclick="markCard('know')">&#10003; Know It</button>
  <button class="score-unsure" onclick="markCard('unsure')">&#126; Unsure</button>
  <button class="score-dunno" onclick="markCard('dunno')">&#10007; Don't Know</button>
</div>
<div class="score-display" id="scoreDisplay"></div>

<div class="keyboard-hint">
  Keyboard: SPACE = flip | LEFT/RIGHT = prev/next | 1 = know | 2 = unsure | 3 = don't know
</div>

<script>
const ALL_CARDS = {cards_json};

let cards = [...ALL_CARDS];
let idx = 0;
let flippedState = false;
let scores = {{ know: 0, unsure: 0, dunno: 0 }};
let activeFilter = "All";

function render() {{
  const c = cards[idx];
  document.getElementById("qCat").textContent = c.cat;
  document.getElementById("qText").textContent = c.q;
  document.getElementById("aCat").textContent = c.cat;
  document.getElementById("aText").textContent = c.a;
  document.getElementById("progress").textContent =
    `Card ${{idx+1}} of ${{cards.length}}  |  Category: ${{c.cat}}`;
  document.getElementById("card").classList.remove("flipped");
  flippedState = false;
  updateScore();
}}

function flip() {{
  document.getElementById("card").classList.toggle("flipped");
  flippedState = !flippedState;
}}

function next() {{
  if (idx < cards.length - 1) idx++;
  else idx = 0;
  render();
}}

function prev() {{
  if (idx > 0) idx--;
  else idx = cards.length - 1;
  render();
}}

function shuffleCards() {{
  for (let i = cards.length - 1; i > 0; i--) {{
    const j = Math.floor(Math.random() * (i + 1));
    [cards[i], cards[j]] = [cards[j], cards[i]];
  }}
  idx = 0;
  render();
}}

function resetCards() {{
  cards = activeFilter === "All"
    ? [...ALL_CARDS]
    : ALL_CARDS.filter(c => c.cat === activeFilter);
  idx = 0;
  scores = {{ know: 0, unsure: 0, dunno: 0 }};
  render();
}}

function markCard(type) {{
  scores[type]++;
  updateScore();
  next();
}}

function updateScore() {{
  const total = scores.know + scores.unsure + scores.dunno;
  document.getElementById("scoreDisplay").textContent =
    `Score: ${{scores.know}} Know | ${{scores.unsure}} Unsure | ${{scores.dunno}} Don't Know | ${{total}} reviewed`;
}}

// Filter buttons
function buildFilters() {{
  const cats = ["All", ...new Set(ALL_CARDS.map(c => c.cat))];
  const bar = document.getElementById("filterBar");
  cats.forEach(cat => {{
    const btn = document.createElement("button");
    btn.className = "filter-btn" + (cat === "All" ? " active" : "");
    btn.textContent = cat;
    btn.onclick = () => {{
      activeFilter = cat;
      cards = cat === "All" ? [...ALL_CARDS] : ALL_CARDS.filter(c => c.cat === cat);
      idx = 0;
      scores = {{ know: 0, unsure: 0, dunno: 0 }};
      document.querySelectorAll(".filter-btn").forEach(b => b.classList.remove("active"));
      btn.classList.add("active");
      render();
    }};
    bar.appendChild(btn);
  }});
}}

// Keyboard
document.addEventListener("keydown", e => {{
  if (e.code === "Space") {{ e.preventDefault(); flip(); }}
  else if (e.code === "ArrowRight") next();
  else if (e.code === "ArrowLeft") prev();
  else if (e.key === "1") markCard("know");
  else if (e.key === "2") markCard("unsure");
  else if (e.key === "3") markCard("dunno");
}});

buildFilters();
render();
</script>
</body>
</html>"""

    with open(HTML_FILE, "w", encoding="utf-8") as f:
        f.write(html)
    print(f"  HTML -> {HTML_FILE}")


# ═════════════════════════════════════════════════════════════════════════════
#  MAIN
# ═════════════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    print(f"\nMontana Motorcycle Test - Flashcard Generator")
    print(f"{'='*48}")
    print(f"Total flashcards: {len(FLASHCARDS)}")
    print(f"Categories: {len(set(c for c,_,_ in FLASHCARDS))}\n")
    print("Generating outputs:")

    generate_csv()
    generate_pdf()
    generate_pptx()
    generate_html()

    print(f"\nAll files saved to: {OUT_DIR}")
    print("Done!")
